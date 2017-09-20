# -*- coding: utf-8 -*-

from getinfo.getTitles import getTitles
from getinfo.getTitles import getCarrers
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import *
from maths import *
import pandas as pd
# Templates of slides
SLD_LAYOUT_ONLY_TITLE_PRESENTATION = 0
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_SECTION_HEADER = 2
SLD_LAYOUT_TWO_SECTION = 3
SLD_LAYOUT_COMPARATION = 4
SLD_LAYOUT_ONLY_TITLE = 5
SLD_LAYOUT_BLANK = 6
SLD_LAYOUT_CONTENT_AND_CAPTION = 7
SLD_LAYOUT_PICTURE_WITH_CAPTION = 8
# This two aren't some common , usually used when work in Japanese lenguage
SLD_LAYOUT_TITLE_AND_VERTICAL_TEXT = 9
SLD_LAYOUT_VERTICAL_TITLE_AND_TEXT = 10

PRS = Presentation()  # Presentation Object , with this object can be created new slides


def generateMultiPPTx():
    """
    Generate the pptx files for all the carrers
    """
    #list_carrer = getCarrers()
    list_carrer = ["Ingenieria Industrial", "Ingenieria Informatica"]
    for carrer in list_carrer:
        generatePPTx(carrer)


def generatePPTx(carrer):
    """
    Generate the pptx for one carrer
    KeysArguments:
        carrer: Name of the carrer
    """
    # idx , for each placeholder is in INF_layouts_placeholder.txt
    list_title = getTitles()
    newSlide = create_slide(PRS, SLD_LAYOUT_ONLY_TITLE_PRESENTATION)
    for title in list_title[1:]:
        create_slide_whith_only_title(PRS, SLD_LAYOUT_ONLY_TITLE, title)

    newSlide = PRS.slides[0]
    newSlide.placeholders[0].text = "Distribución de la Demanda Laboral".upper(
    )
    newSlide.placeholders[1].text = carrer.title()
    PRS.save(carrer + ".pptx")


def create_slide(prs, constSLDLayout,title=None):
    """Create any type of slide 


    Arguments:
        PRS {Presentation-Object} 
        constSLDLayout {Constant for the type of slide}

    Returns:
        Slide -- The new slide added 
    """
  

    slide = prs.slide_layouts[constSLDLayout]
    newSlide = prs.slides.add_slide(slide)
    if title is not None :
        newSlide.placeholders[0].text = title
    return newSlide


def create_slide_whith_only_title(prs, constSLDLayout, title=None):
    """
    Create a new Slide
    KeysArguments
    PRS = Presentation object
    constSLDLayou = Type of layout
    title = Title of the Slide , can be omitted
    """
    slide = prs.slide_layouts[constSLDLayout]
    newSlide = prs.slides.add_slide(slide)
    if title is not None:
        newSlide.placeholders[0].text = title
    return newSlide


def create_slide_with_one_image(prs, constSLDLayout, path=None, tuplePosition=None, tupleSize=None):
    """
    Create a new slide with a image
    KeysArguments
    PRS = Presentation object
    path = Path of the image
    constSLDLayout = Type of layout
    tuplePosition(top, left) = Tuple of the position on the top-left point of the image
    tupleSize(width, ) = Tuple of the size (width , height)
    """
    assert path is not None, "Can't fine the path of the image"
    slide = PRS.slide_layouts[constSLDLayout]
    newSlide = PRS.slides.add_slide(slide)
    assert tuplePosition is not None, "The image need a position "
    left = tuplePosition[0]
    top = tuplePosition[1]
    if (tupleSize is not None):
        width = tupleSize[0]
        height = tupleSize[1]
    newSlide.add_image(path, left, top, width, height)


def add_chart_to_slide(slide,info, categories, tuplePosition=None):
    assert path is not None, "The chart need a position"
    x = tuplePosition[0]
    y = tuplePosition[1]
    char_data = ChartData()
    char_data.categories = categories
    char_data.add_series('Series 1', info)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, char_data).chart
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = True


def add_table_to_slide(slide, rows, cols):
    shapes = slide.shapes
    width = Inches(2.0)
    height = Inches(0.8)
    left = top = Inches(2.0)
    table = shapes.add_table(rows + 1, cols + 1, left,
                             top, width, height).table
    for columns in table.columns:
        columns.width = Inches(2.75)
    return table


def export_ppt(list_reports, list_file):
    assert len(list_reports) == len(
        list_file), "The list dosent have the same size"
    for index, _type in enumarate(list_reports):
        title = list_file[index].split(".")[0].title()
        newSlide = create_slide_whith_only_title(
            PRS, SLD_LAYOUT_ONLY_TITLE, title)
        if _type == "pie":
            dateFrame = pd.read_csv(list_file[index])
            list_unique_values = list(
                dateFrame[title].unique())
            value = dateFrame[title].sum()
            info = []
            for x in dateFrame[title]:
                info.append(float(x)/value)
            add_chart_to_slide(newSlide,info,list_unique_values,tuplePosition=(Inches(2.0),Inches(2.0)))
        elif _type == "bar":
            pass
        elif _type == "table_float":
            matrix = calculate_alpha(list_file[index])
            list_unique_values = list(
                pd.read_csv(list_file[index]).Area.unique())
            lenAreas = len(list_unique_values)
            table = add_table_to_slide(newSlide, rows=lenAreas, cols=lenAreas)
            for indexX, _ in enumerate(lenAreas):
                for indexY, _ in enumerate(lenAreas):
                    table.cell(indexX + 1, indexY +1).text = str(matrix[indexX, indexY] * 100) + "%"
        elif _type == "list_table":
            pass 
             #Global pd
