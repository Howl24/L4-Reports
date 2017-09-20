# -*- coding: utf-8 -*-
from pptx import Presentation
import sys
def getTitles():
    """Get all the titles from a Pptx template """
    titles = list()
    try:
        prs = Presentation("Template.pptx")
        file = open("output.txt","w")
        sys.output = file
        for slide in prs.slides:
            line = "".join(slide.placeholders[0].text)
            titles.append(line+"\n")
    except FileNotFoundError :
        pass
    return titles



def getCarrers():
    """Get all the careers """
    careers = []
    file = open("Careers.txt","r")
    for career in file:
        careers.append(career)
    return careers
