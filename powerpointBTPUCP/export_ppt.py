
# coding: utf-8

from .PandasToPowerpoint import *
import pandas as pd
from pptx import Presentation
from ppxtools.supportPP import *
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import *
from datetime import datetime as dt

dictionryTitle = {"laboral":"Demanda Laboral","cargo":"Según Cargo","nivelEstudio":"Según nivel de estudios demandados","conocimiento":"Conocimientos adicionales","profesional":"Perfil Profesional","cruces":"Perfil Profesional: Cruces","idiomas":"Idiomas demandados"}



def appendDF(filename_list,key):
    df = None
    for filename in filename_list:
        name = filename.split("-")
        name = name[1][:]
        keyInFile = name.split(".")[0]
        if keyInFile == key :
            if df is None:
                df = pd.read_csv(filename)
            else:
                df.append(pd.read_csv(filename))
        return df

def create_post_titleSlide(prs,filename_list):
    list_titles = ["Según tamaño de empresa","Según sectores de actividad","Según cargo","Según nivel de estudios demandados","Conocimientos demandados","Según perfiles profesionales de la universidad"]
    for title in list_titles:
        newSlide = create_slide(prs,SLD_LAYOUT_ONLY_TITLE,title)
        if title == list_titles[2]:
            df = appendDF(filename_list,"cargo")
            if df is not None :
                draw_char(newSlide, "" , df=df)
            #Example

def createBasicTitleSlide(prs,carrer):
    slide = prs.slide_layouts[SLD_LAYOUT_ONLY_TITLE_PRESENTATION]
    newSlide = prs.slides.add_slide(slide)
    newSlide.placeholders[0].text = "Distribución de Demanda Laboral".title()
    newSlide.placeholders[1].text = carrer.upper() + "-" + str(dt.now().year)

def createBasicTitleSlideForService(prs,serviceName):
    slide = prs.slide_layouts[SLD_LAYOUT_ONLY_TITLE_PRESENTATION]
    newSlide = prs.slides.add_slide(slide)
    newSlide.placeholders[0].text = "Sistema de Oportunidades Laborales".title()
    newSlide.placeholders[1].text = serviceName.title()

def createBasicPresenSlide(prs,title):
    slide = prs.slide_layouts[SLD_LAYOUT_ONLY_TITLE]
    newSlide = prs.slides.add_slide(slide)
    newSlide.placeholders[0].text = dictionryTitle[title]
    return newSlide

def draw_char(slide,filename_list,df=None):
    x = y = Inches(1.5)
    cx = Inches(8)
    cy = Inches(6)
    char_data = ChartData()
    if df is None:
        info = calculate_percentage(nameFile=filename_list)
        data = pd.read_csv(filename_list)
    else :
        data = df
        info = calculate_percentage(nameFile = "",df=df)
    #info[-1] += (1-info[-1])
    name = data.columns[0]
    char_data.categories = [ x.upper() for x in list(data[name].unique())]
    char_data.add_series('Series 1', info)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, char_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.LEFT
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END



def draw_bar(slide,filename_list,df=None):
    chart_data = ChartData()
    if df is None:
        data = pd.read_csv(filename_list)
    else :
        data = df
    name=data.columns[0]
    chart_data.categories = [ x.upper() for x in list(data[name].unique())]
    name = data.columns[1]
    chart_data.add_series('Series 1', data[name])
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

def add_tablePer(slide,filename_list,df=None):
    if df is None:
        data = pd.read_csv(filename_list)
    else :
        data = df  
    df_to_table(slide,df,left=Inches(2),top=Inches(2),width=Inches(4),heigh=Inches(6))
def add_tableList(slide,filename_list,df=None):
    if df is None:
        data = pd.read_csv(filename_list)
    else :
        data = df  
    name = data.columns[0]
    colnames = list(data[name].unique())
    df_to_table(slide,df,left=Inches(2),top=Inches(2),width=Inches(4),heigh=Inches(6),colnames=colnames,rowname=colnames)

def add_tableListNumber(slide,filename_list):
    #Finalice this 
    pass


def calculate_percentage(nameFile,df=None):
    if df is None :
        data = pd.read_csv(nameFile)
    else :
        data = df
    sumValues = values = data.sum().values[-1]
    listDoubles = len(data)*[0]
    for x,y in enumerate(listDoubles):
        listDoubles[x] = float(data.iloc[x,1])/sumValues
    return listDoubles

dictionary_functions = {"pie":draw_char,"bar":draw_bar,"per_table":add_tablePer,"list_table":add_tableList,"list_number_str":add_tableListNumber}
def export_ppt(report_type_list,filename_list,carrer="ECONOMÍA"):
    assert len(report_type_list) == len(filename_list), "The size of the two list is not the size "
    prs = Presentation()
    createBasicTitleSlide(prs,carrer)
    #NeewFunction to create previusSlides
    create_post_titleSlide(prs,filename_list[:])
    webPages = filename_list[0].split('-')[:]
    createBasicTitleSlideForService(prs,webPages[0])
    pastName= webPages[0]
    for type_ ,filename in list(zip(report_type_list,filename_list[:])):
        webPages = filename.split('-')[:]
        if webPages[0] != pastName :
            createBasicTitleSlideForService(prs,webPages[0])
        else:
            title = webPages[1].split(".")[0]
            slide = createBasicPresenSlide(prs,title)
            try:
                dictionary_functions[type_](slide,filename)
            except KeyError:
                continue
    prs.save(carrer + ".pptx")
